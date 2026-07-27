import os
import sys
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Define Color Palette
C_PRIMARY = RGBColor(24, 76, 52)      # Deep Forest Green
C_ACCENT = RGBColor(138, 171, 143)     # Sage Green
C_BG_LIGHT = RGBColor(248, 249, 250)   # Off-White / Light Pearl
C_TEXT_DARK = RGBColor(33, 37, 41)     # Charcoal / Near Black
C_TEXT_LIGHT = RGBColor(255, 255, 255) # White
C_CARD_BG = RGBColor(238, 244, 240)    # Soft Green-Grey Card BG
C_CARD_BORDER = RGBColor(210, 225, 215)# Card Border

def create_slide_title(slide, prs, title_text, subtitle_text=None):
    """Creates a unified header with Title and optional Subtitle."""
    # Text box for Title & Subtitle
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Title Paragraph
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = C_PRIMARY
    
    # Subtitle Paragraph (if provided)
    if subtitle_text:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = C_ACCENT
        p_sub.space_before = Pt(4)
        
    # Decorative line under the header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

def add_bullet_point(tf, bold_prefix, normal_text, font_size=13, is_dark_bg=False):
    """Adds a bullet point with bold prefix and normal description."""
    p = tf.add_paragraph()
    p.space_after = Pt(6)
    p.level = 0
    
    # Colors based on background
    color_bold = C_TEXT_LIGHT if is_dark_bg else C_PRIMARY
    color_normal = C_TEXT_LIGHT if is_dark_bg else C_TEXT_DARK
    
    # Bullet character
    r_bullet = p.add_run()
    r_bullet.text = "• "
    r_bullet.font.name = "Calibri"
    r_bullet.font.size = Pt(font_size)
    r_bullet.font.color.rgb = color_bold
    
    # Bold prefix
    r_bold = p.add_run()
    r_bold.text = bold_prefix
    r_bold.font.name = "Calibri"
    r_bold.font.bold = True
    r_bold.font.size = Pt(font_size)
    r_bold.font.color.rgb = color_bold
    
    # Normal text
    r_normal = p.add_run()
    r_normal.text = normal_text
    r_normal.font.name = "Calibri"
    r_normal.font.size = Pt(font_size)
    r_normal.font.color.rgb = color_normal

def create_card_shape(slide, left, top, width, height, fill_color, border_color):
    """Creates a beautifully styled rounded rectangle card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def build_presentation(output_path):
    print("Initializing Presentation...")
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme for Visual Impact)
    # ----------------------------------------------------
    print("Building Slide 1: Title...")
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Dark Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_PRIMARY
    bg1.line.fill.background()
    
    # Decorative Accent Box
    accent_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), prs.slide_height)
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = C_ACCENT
    accent_box.line.fill.background()
    
    # Title Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "AeroCampus-AI"
    p1.font.name = "Arial"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = C_TEXT_LIGHT
    p1.space_after = Pt(10)
    
    p2 = tf1.add_paragraph()
    p2.text = "Smart Campus Environmental Management & Mitigation System"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_ACCENT
    p2.space_after = Pt(40)
    
    p3 = tf1.add_paragraph()
    p3.text = "Assignment 4B: Solution Scaling & Implementation Strategy"
    p3.font.name = "Calibri"
    p3.font.size = Pt(16)
    p3.font.color.rgb = C_TEXT_LIGHT
    
    p4 = tf1.add_paragraph()
    p4.text = "Author: Sahil Agarwal  |  Context: JIIT Noida Sector 62 Campus"
    p4.font.name = "Calibri"
    p4.font.size = Pt(14)
    p4.font.color.rgb = C_ACCENT
    
    # ----------------------------------------------------
    # SLIDE 2: Target Users & Value Proposition
    # ----------------------------------------------------
    print("Building Slide 2: Target Users & Value Proposition...")
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Light Background
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_BG_LIGHT
    bg2.line.fill.background()
    
    create_slide_title(slide2, prs, "Target Users & Value Proposition", "Aligning health benefits with administrative operations")
    
    # Left Card: Target Users
    create_card_shape(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), C_CARD_BG, C_CARD_BORDER)
    left_tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
    tf_left = left_tb.text_frame
    tf_left.word_wrap = True
    
    p_left_title = tf_left.paragraphs[0]
    p_left_title.text = "TARGET USERS & PROBLEMS"
    p_left_title.font.name = "Arial"
    p_left_title.font.size = Pt(18)
    p_left_title.font.bold = True
    p_left_title.font.color.rgb = C_PRIMARY
    p_left_title.space_after = Pt(14)
    
    add_bullet_point(tf_left, "Students & Faculty: ", "Noida Sector 62 campus residents exposed to raw highway PM2.5 levels exceeding 140 ug/m3 in winter sports and classroom settings.")
    add_bullet_point(tf_left, "Sensitive Cohorts: ", "Asthmatic and allergy-prone members needing instant, micro-local warnings to preempt acute health episodes.")
    add_bullet_point(tf_left, "Campus Operations: ", "Security and facilities teams requiring immediate decision rules to restrict gate traffic and optimize HVAC system configurations.")
    
    # Right Card: Value Proposition (Dark Forest Green for Contrast)
    create_card_shape(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), C_PRIMARY, None)
    right_tb = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.6))
    tf_right = right_tb.text_frame
    tf_right.word_wrap = True
    
    p_right_title = tf_right.paragraphs[0]
    p_right_title.text = "VALUE PROPOSITION & IMPACT"
    p_right_title.font.name = "Arial"
    p_right_title.font.size = Pt(18)
    p_right_title.font.bold = True
    p_right_title.font.color.rgb = C_ACCENT
    p_right_title.space_after = Pt(14)
    
    add_bullet_point(tf_right, "Exposure Mitigation (>95%): ", "Proactive push alerts to sports coordinators reduce student exposure hours to hazardous PM2.5 levels on campus to near zero.", is_dark_bg=True)
    add_bullet_point(tf_right, "Indoor Air Quality (40%-60%): ", "Closed-loop automation pre-conditions buildings before traffic peaks occur, maintaining clean baselines.", is_dark_bg=True)
    add_bullet_point(tf_right, "Operational Efficiency: ", "Automates physical security gates and ventilation dampers based on AI triggers, eliminating manual errors.", is_dark_bg=True)
    
    # ----------------------------------------------------
    # SLIDE 3: Scaling Strategy
    # ----------------------------------------------------
    print("Building Slide 3: Scaling Strategy...")
    slide3 = prs.slides.add_slide(blank_layout)
    
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = C_BG_LIGHT
    bg3.line.fill.background()
    
    create_slide_title(slide3, prs, "Solution Scaling Strategy", "Multi-dimensional growth paths across locations, user groups, and partners")
    
    col_width = Inches(3.6)
    col_height = Inches(5.0)
    col_y = Inches(1.8)
    
    # Column 1: Across Locations
    create_card_shape(slide3, Inches(0.8), col_y, col_width, col_height, C_CARD_BG, C_CARD_BORDER)
    tb_c1 = slide3.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(3.4), Inches(4.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    
    p_c1_h = tf_c1.paragraphs[0]
    p_c1_h.text = "1. ACROSS LOCATIONS"
    p_c1_h.font.name = "Arial"
    p_c1_h.font.size = Pt(16)
    p_c1_h.font.bold = True
    p_c1_h.font.color.rgb = C_PRIMARY
    p_c1_h.space_after = Pt(12)
    
    add_bullet_point(tf_c1, "Sister Institutions: ", "Expand from JIIT Noida Sector 62 to other Jaypee campuses (Sector 128, Waknaghat, Guna).")
    add_bullet_point(tf_c1, "AeroCampus-in-a-Box: ", "Package the AI-IoT sensor suite and gateway code for easy installation at other high-density academic and corporate campuses in Delhi-NCR.")
    
    # Column 2: Across User Groups
    create_card_shape(slide3, Inches(4.8), col_y, col_width, col_height, C_CARD_BG, C_CARD_BORDER)
    tb_c2 = slide3.shapes.add_textbox(Inches(4.9), Inches(1.9), Inches(3.4), Inches(4.8))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    
    p_c2_h = tf_c2.paragraphs[0]
    p_c2_h.text = "2. ACROSS USER GROUPS"
    p_c2_h.font.name = "Arial"
    p_c2_h.font.size = Pt(16)
    p_c2_h.font.bold = True
    p_c2_h.font.color.rgb = C_PRIMARY
    p_c2_h.space_after = Pt(12)
    
    add_bullet_point(tf_c2, "Surrounding Noida RWAs: ", "Extend notifications and data access to neighboring residential welfare associations and public schools.")
    add_bullet_point(tf_c2, "Personal App Features: ", "Introduce personalized student exposure logs and optimized 'clean air routes' mapped across campus corridors.")
    
    # Column 3: Through Partnerships
    create_card_shape(slide3, Inches(8.8), col_y, col_width, col_height, C_CARD_BG, C_CARD_BORDER)
    tb_c3 = slide3.shapes.add_textbox(Inches(8.9), Inches(1.9), Inches(3.4), Inches(4.8))
    tf_c3 = tb_c3.text_frame
    tf_c3.word_wrap = True
    
    p_c3_h = tf_c3.paragraphs[0]
    p_c3_h.text = "3. PARTNERSHIPS"
    p_c3_h.font.name = "Arial"
    p_c3_h.font.size = Pt(16)
    p_c3_h.font.bold = True
    p_c3_h.font.color.rgb = C_PRIMARY
    p_c3_h.space_after = Pt(12)
    
    add_bullet_point(tf_c3, "Smart City Integration: ", "Share real-time micro-local pollution maps with the Noida Municipal Authority's public dashboards.")
    add_bullet_point(tf_c3, "Horticultural Alliances: ", "Partner with agricultural institutes (e.g., IARI) to audit and refine species-specific leaf deposition parameters ($k$) for urban forest barriers.")
    
    # ----------------------------------------------------
    # SLIDE 4: Phased Implementation Plan
    # ----------------------------------------------------
    print("Building Slide 4: Implementation Plan...")
    slide4 = prs.slides.add_slide(blank_layout)
    
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = C_BG_LIGHT
    bg4.line.fill.background()
    
    create_slide_title(slide4, prs, "Phased Implementation Plan", "Staged rollout to validate ML models, automate HVAC, and establish the Green Shield")
    
    card_w = Inches(3.6)
    card_h = Inches(4.2)
    card_y = Inches(2.2)
    
    # Draw horizontal timeline arrow
    tl_line = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.2))
    tl_line.fill.solid()
    tl_line.fill.fore_color.rgb = C_ACCENT
    tl_line.line.fill.background()
    
    # Phase 1: Pilot
    create_card_shape(slide4, Inches(0.8), card_y, card_w, card_h, C_CARD_BG, C_CARD_BORDER)
    tb_p1 = slide4.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(3.4), Inches(4.0))
    tf_p1 = tb_p1.text_frame
    tf_p1.word_wrap = True
    
    p_p1_h = tf_p1.paragraphs[0]
    p_p1_h.text = "PHASE 1: PILOT (M1-M4)"
    p_p1_h.font.name = "Arial"
    p_p1_h.font.size = Pt(15)
    p_p1_h.font.bold = True
    p_p1_h.font.color.rgb = C_PRIMARY
    p_p1_h.space_after = Pt(10)
    
    add_bullet_point(tf_p1, "Setup: ", "Deploy 5 outdoor and 5 indoor IoT sensor pods. Establish data pipeline.")
    add_bullet_point(tf_p1, "Validation: ", "Calibrate Random Forest ML model on local wind/traffic proxies.")
    add_bullet_point(tf_p1, "Next-Phase Trigger: ", "$>85\%$ prediction accuracy and zero packet loss for 30 consecutive days.")
    
    # Phase 2: Expansion
    create_card_shape(slide4, Inches(4.8), card_y, card_w, card_h, C_CARD_BG, C_CARD_BORDER)
    tb_p2 = slide4.shapes.add_textbox(Inches(4.9), Inches(2.3), Inches(3.4), Inches(4.0))
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True
    
    p_p2_h = tf_p2.paragraphs[0]
    p_p2_h.text = "PHASE 2: EXPANSION (M5-M9)"
    p_p2_h.font.name = "Arial"
    p_p2_h.font.size = Pt(15)
    p_p2_h.font.bold = True
    p_p2_h.font.color.rgb = C_PRIMARY
    p_p2_h.space_after = Pt(10)
    
    add_bullet_point(tf_p2, "Setup: ", "Automate indoor air purifiers and HVAC dampers. Roll out student alert app.")
    add_bullet_point(tf_p2, "Policy: ", "Execute Gate Traffic restrictions during peak smog days.")
    add_bullet_point(tf_p2, "Next-Phase Trigger: ", "Compliance rate (G1) $>90\%$ and classroom PM2.5 kept below 35 ug/m3.")
    
    # Phase 3: Full-Scale
    create_card_shape(slide4, Inches(8.8), card_y, card_w, card_h, C_CARD_BG, C_CARD_BORDER)
    tb_p3 = slide4.shapes.add_textbox(Inches(8.9), Inches(2.3), Inches(3.4), Inches(4.0))
    tf_p3 = tb_p3.text_frame
    tf_p3.word_wrap = True
    
    p_p3_h = tf_p3.paragraphs[0]
    p_p3_h.text = "PHASE 3: FULL-SCALE (M10-M12)"
    p_p3_h.font.name = "Arial"
    p_p3_h.font.size = Pt(15)
    p_p3_h.font.bold = True
    p_p3_h.font.color.rgb = C_PRIMARY
    p_p3_h.space_after = Pt(10)
    
    add_bullet_point(tf_p3, "Green Shield: ", "Plant 1,500 m2 dense Neem and Pilkan canopy on the NE campus perimeter.")
    add_bullet_point(tf_p3, "Automation: ", "Link physical motorized gate barriers directly to AI policy outputs.")
    add_bullet_point(tf_p3, "Audit: ", "Conduct annual satellite canopy cover assessments and monthly ESG reporting.")
    
    # ----------------------------------------------------
    # SLIDE 5: Risks & Mitigation
    # ----------------------------------------------------
    print("Building Slide 5: Risks & Mitigation...")
    slide5 = prs.slides.add_slide(blank_layout)
    
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = C_BG_LIGHT
    bg5.line.fill.background()
    
    create_slide_title(slide5, prs, "Risk Management & Mitigation Framework", "Addressing operational, technical, and biological hurdles")
    
    row_w = Inches(11.7)
    row_h = Inches(1.5)
    row_x = Inches(0.8)
    
    # Risk 1 Row
    y1 = Inches(1.8)
    create_card_shape(slide5, row_x, y1, row_w, row_h, C_CARD_BG, C_CARD_BORDER)
    tb_r1 = slide5.shapes.add_textbox(row_x + Inches(0.2), y1 + Inches(0.1), row_w - Inches(0.4), row_h - Inches(0.2))
    tf_r1 = tb_r1.text_frame
    tf_r1.word_wrap = True
    
    p_r1 = tf_r1.paragraphs[0]
    p_r1.text = "TECHNICAL: Sensor Drift and Soot Accumulation"
    p_r1.font.name = "Arial"
    p_r1.font.size = Pt(14)
    p_r1.font.bold = True
    p_r1.font.color.rgb = C_PRIMARY
    p_r1.space_after = Pt(6)
    
    add_bullet_point(tf_r1, "Impact: ", "Decaying accuracy causes false alarms or delayed HVAC activations.")
    add_bullet_point(tf_r1, "Mitigation: ", "Deploy daily auto-calibration routines using regional CPCB stations; quarterly physical filter/sensor cleaning.")
    
    # Risk 2 Row
    y2 = Inches(3.5)
    create_card_shape(slide5, row_x, y2, row_w, row_h, C_CARD_BG, C_CARD_BORDER)
    tb_r2 = slide5.shapes.add_textbox(row_x + Inches(0.2), y2 + Inches(0.1), row_w - Inches(0.4), row_h - Inches(0.2))
    tf_r2 = tb_r2.text_frame
    tf_r2.word_wrap = True
    
    p_r2 = tf_r2.paragraphs[0]
    p_r2.text = "OPERATIONAL: Low User Compliance & Alert Fatigue"
    p_r2.font.name = "Arial"
    p_r2.font.size = Pt(14)
    p_r2.font.bold = True
    p_r2.font.color.rgb = C_PRIMARY
    p_r2.space_after = Pt(6)
    
    add_bullet_point(tf_r2, "Impact: ", "Sports coordinators ignore outdoor suspension alerts; gate guards fail to check non-compliant vehicles.")
    add_bullet_point(tf_r2, "Mitigation: ", "Link the physical gate barrier directly to the AI simulator; tie alert response times to departmental performance logs.")
    
    # Risk 3 Row
    y3 = Inches(5.2)
    create_card_shape(slide5, row_x, y3, row_w, row_h, C_CARD_BG, C_CARD_BORDER)
    tb_r3 = slide5.shapes.add_textbox(row_x + Inches(0.2), y3 + Inches(0.1), row_w - Inches(0.4), row_h - Inches(0.2))
    tf_r3 = tb_r3.text_frame
    tf_r3.word_wrap = True
    
    p_r3 = tf_r3.paragraphs[0]
    p_r3.text = "ECOLOGICAL: Vegetative Shield Growth Lag"
    p_r3.font.name = "Arial"
    p_r3.font.size = Pt(14)
    p_r3.font.bold = True
    p_r3.font.color.rgb = C_PRIMARY
    p_r3.space_after = Pt(6)
    
    add_bullet_point(tf_r3, "Impact: ", "Saplings require 3 to 5 years to form a dense canopy, leaving the highway boundary unshielded in the short term.")
    add_bullet_point(tf_r3, "Mitigation: ", "Install dual-layer synthetic dust filters on the border fence; intersperse Neem with fast-growing hybrid Poplars for early cover.")
    
    # Save Presentation
    print(f"Saving presentation to {output_path}...")
    prs.save(output_path)
    print("Presentation saved successfully.")

def convert_to_pdf(pptx_path, pdf_path):
    print("Attempting to convert PPTX to PDF via PowerPoint COM Automation...")
    try:
        import win32com.client
        import pythoncom
        
        # Initialize COM
        pythoncom.CoInitialize()
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # Keep presentation window hidden
        pres = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        
        print(f"Saving PDF to: {pdf_path}")
        pres.SaveAs(os.path.abspath(pdf_path), FileFormat=32) # 32 is ppSaveAsPDF
        pres.Close()
        powerpoint.Quit()
        print("PDF conversion completed successfully.")
        return True
    except Exception as e:
        print(f"Could not convert to PDF using PowerPoint COM Automation: {e}")
        print("Make sure Microsoft PowerPoint is installed and COM is working.")
        return False

def main():
    workspace_dir = r"d:\1m1b internship\project"
    pptx_filename = "SahilAgarwal_ScalingStrategy.pptx"
    pdf_filename = "SahilAgarwal_ScalingStrategy.pdf"
    
    pptx_path = os.path.join(workspace_dir, pptx_filename)
    pdf_path = os.path.join(workspace_dir, pdf_filename)
    
    build_presentation(pptx_path)
    convert_to_pdf(pptx_path, pdf_path)

if __name__ == "__main__":
    main()
